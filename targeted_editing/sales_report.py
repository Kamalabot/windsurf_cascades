import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt

# Read Excel data
df = pd.read_excel('SaleData.xlsx')

# Create visualizations
plt.figure(figsize=(10,6))
region_sales = df.groupby('Region')['Sale_amt'].sum()
region_sales.plot(kind='bar')
plt.title('Total Sales by Region')
plt.xlabel('Region')
plt.ylabel('Total Sales ($)')
plt.tight_layout()
plt.savefig('region_sales_chart.png')
plt.close()

plt.figure(figsize=(10,6))
salesman_sales = df.groupby('SalesMan')['Sale_amt'].sum()
salesman_sales.plot(kind='bar')
plt.title('Total Sales by Salesman')
plt.xlabel('Salesman')
plt.ylabel('Total Sales ($)')
plt.tight_layout()
plt.savefig('salesman_sales_chart.png')
plt.close()

# Create Presentation
prs = Presentation()

# Slide 1: Overview
slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = 'Sales Performance Overview'
subtitle.text = f"""
Total Sales: ${df['Sale_amt'].sum():,.2f}
Average Sale: ${df['Sale_amt'].mean():,.2f}
Total Units Sold: {df['Units'].sum()}
Date Range: {df['OrderDate'].min().date()} to {df['OrderDate'].max().date()}
"""

# Slide 2: Regional Sales Analysis
slide2 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide

# Add title manually
left = Inches(1)
top = Inches(0.5)
width = Inches(8)
height = Inches(1)
txBox = slide2.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = 'Sales by Region'
p.font.size = Pt(24)
p.font.bold = True

# Add chart to slide
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(5)
pic = slide2.shapes.add_picture('region_sales_chart.png', left, top, width, height)

# Add text box with regional sales summary
left = Inches(1)
top = Inches(6.5)
width = Inches(8)
height = Inches(1)
txBox = slide2.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
region_sales = df.groupby('Region')['Sale_amt'].sum().sort_values(ascending=False)
p.text = "Regional Sales Summary: " + ", ".join([f"{region}: ${sales:,.2f}" for region, sales in region_sales.items()])

# Slide 3: Salesman Sales Analysis
slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide

# Add title manually
left = Inches(1)
top = Inches(0.5)
width = Inches(8)
height = Inches(1)
txBox = slide3.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = 'Sales by Salesman'
p.font.size = Pt(24)
p.font.bold = True

# Add chart to slide
left = Inches(1)
top = Inches(1.5)
width = Inches(8)
height = Inches(5)
pic = slide3.shapes.add_picture('salesman_sales_chart.png', left, top, width, height)

# Add text box with salesman sales summary
left = Inches(1)
top = Inches(6.5)
width = Inches(8)
height = Inches(1)
txBox = slide3.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
salesman_sales = df.groupby('SalesMan')['Sale_amt'].sum().sort_values(ascending=False)
p.text = "Salesman Sales Summary: " + ", ".join([f"{salesman}: ${sales:,.2f}" for salesman, sales in salesman_sales.items()])

# Save Presentation
prs.save('Sales_Report.pptx')
print("Presentation created successfully with charts!")
